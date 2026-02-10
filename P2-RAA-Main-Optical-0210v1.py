import streamlit as st
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import io

# 設定頁面標題
st.set_page_config(page_title="工廠數據報告產生器", page_icon="📊")

st.title("🏭 工廠數據自動化報告產生器")
st.markdown("### 請上傳 Excel 原始數據 (RAA / IPQC)")

# --- 核心處理函數 ---

def find_header_row(file_obj, sheet_name):
    try:
        file_obj.seek(0) # 重置指標
        df_temp = pd.read_excel(file_obj, sheet_name=sheet_name, header=None, nrows=20, engine='openpyxl')
        for idx, row in df_temp.iterrows():
            if isinstance(row[0], str) and row[0].strip().startswith('Tester'):
                return idx
        return 0
    except:
        return 0

def get_station_name(col_name):
    if 'PreAA' in col_name:
        if 'H1' in col_name or 'V1' in col_name: return 'PreAA_1'
        if 'H2' in col_name or 'V2' in col_name: return 'PreAA_2'
    if 'AfterExposure' in col_name: return 'AfterExp'
    if 'LooseClaws' in col_name: return 'LooseClaws'
    if 'AA_M87' in col_name: return 'AA'
    if 'AfterBaking' in col_name: return 'AfterBaking'
    return None

def process_uploaded_files(uploaded_files):
    all_data = []
    
    for uploaded_file in uploaded_files:
        try:
            xls = pd.ExcelFile(uploaded_file, engine='openpyxl')
            for sheet in xls.sheet_names:
                if sheet not in ['RAA-R', 'RAA-L', 'IPQC-R', 'IPQC-L']:
                    continue
                
                header_idx = find_header_row(uploaded_file, sheet)
                uploaded_file.seek(0) # 重要：讀取前要重置檔案指標
                df = pd.read_excel(uploaded_file, sheet_name=sheet, header=header_idx, engine='openpyxl')
                
                side = 'Right' if '-R' in sheet else 'Left'
                target_cols = [c for c in df.columns if 'Boresight' in str(c) and 'White' in str(c)]
                
                if not target_cols: continue

                if 'CreateTime' in df.columns:
                    df['CreateTime'] = pd.to_datetime(df['CreateTime'], errors='coerce')
                
                melted = df.melt(id_vars=['CreateTime'], value_vars=target_cols, 
                                 var_name='Station_Raw', value_name='Value')
                melted['Side'] = side
                melted['Source'] = uploaded_file.name
                
                def get_direction(name):
                    if '_H_' in name or 'illu_Boresight_H' in name: return 'H'
                    if '_V_' in name or 'illu_Boresight_V' in name: return 'V'
                    return 'Unknown'
                
                melted['Direction'] = melted['Station_Raw'].apply(get_direction)
                melted['Station_Generic'] = melted['Station_Raw'].apply(get_station_name)
                all_data.append(melted)
        except Exception as e:
            st.error(f"讀取檔案 {uploaded_file.name} 失敗: {e}")

    if not all_data: return pd.DataFrame()
    
    final_df = pd.concat(all_data, ignore_index=True)
    final_df['Value'] = pd.to_numeric(final_df['Value'], errors='coerce')
    return final_df.dropna(subset=['Value', 'Station_Generic'])

def generate_ppt(df):
    prs = Presentation()
    sns.set_style("whitegrid")
    station_order = ['PreAA_1', 'PreAA_2', 'AA', 'AfterExp', 'LooseClaws', 'AfterBaking']
    
    # 建立排序標籤
    df['Display_Label'] = df['Side'].str[0] + '-' + df['Station_Generic']
    plot_order = [f"L-{s}" for s in station_order] + [f"R-{s}" for s in station_order]

    # 繪圖暫存區
    image_buffers = {}

    # 1. Box Plot Logic
    def create_boxplot(data, title, ylim=None):
        fig, axes = plt.subplots(2, 1, figsize=(10, 8), sharex=True)
        for i, direction in enumerate(['H', 'V']):
            sub_data = data[data['Direction'] == direction]
            sns.boxplot(data=sub_data, x='Display_Label', y='Value', hue='Side', 
                        order=[x for x in plot_order if x in sub_data['Display_Label'].unique()],
                        palette={'Left':'skyblue', 'Right':'orange'}, ax=axes[i], dodge=False)
            axes[i].set_title(f'{title} - {direction}')
            axes[i].axhline(0.25, color='red', linestyle='--', label='USL')
            axes[i].axhline(-0.25, color='red', linestyle='--', label='LSL')
            if ylim: axes[i].set_ylim(ylim)
            axes[i].grid(True, linestyle=':', alpha=0.6)
            if len(axes[i].get_xticks()) > 0:
                mid = len(axes[i].get_xticks()) / 2 - 0.5
                axes[i].axvline(mid, color='grey', linestyle='-.')
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=100)
        plt.close()
        buf.seek(0)
        return buf

    # 2. Control Chart Logic
    def create_control_chart(data, ylim=None):
        ab_data = data[data['Station_Generic'] == 'AfterBaking'].sort_values('CreateTime')
        fig, axes = plt.subplots(2, 2, figsize=(12, 8))
        colors = {'Left': 'blue', 'Right': 'orange'}
        
        for i, d in enumerate(['H', 'V']):
            for j, s in enumerate(['Left', 'Right']):
                ax = axes[i, j]
                subset = ab_data[(ab_data['Direction'] == d) & (ab_data['Side'] == s)]
                if not subset.empty:
                    ax.scatter(subset['CreateTime'], subset['Value'], color=colors[s], alpha=0.6)
                    plt.setp(ax.get_xticklabels(), rotation=30, ha='right')
                ax.set_title(f'{d} - {s}')
                ax.axhline(0.25, color='red', linestyle='--')
                ax.axhline(-0.25, color='red', linestyle='--')
                ax.grid(True)
                if ylim: ax.set_ylim(ylim)
        
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=100)
        plt.close()
        buf.seek(0)
        return buf

    # 生成圖片
    # Overall
    image_buffers['overall_auto'] = create_boxplot(df, 'Overall Summary')
    image_buffers['overall_fixed'] = create_boxplot(df, 'Overall Summary', (-1.5, 1.5))
    
    # Latest
    latest_date = df['CreateTime'].max().date()
    latest_df = df[df['CreateTime'].dt.date == latest_date].copy()
    image_buffers['latest_auto'] = create_boxplot(latest_df, f'Latest Data ({latest_date})')
    image_buffers['latest_fixed'] = create_boxplot(latest_df, f'Latest Data ({latest_date})', (-1.5, 1.5))

    # Control Chart
    image_buffers['control_auto'] = create_control_chart(df)
    image_buffers['control_fixed'] = create_control_chart(df, (-0.3, 0.3))

    # 製作 PPT
    def add_dual_slide(title, img_key1, img_key2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1)).text_frame.text = title
        slide.shapes.add_picture(image_buffers[img_key1], Inches(0.2), Inches(1.5), width=Inches(4.8))
        slide.shapes.add_picture(image_buffers[img_key2], Inches(5.1), Inches(1.5), width=Inches(4.8))
        # Add labels
        slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(2), Inches(0.5)).text_frame.text = "Auto Scale"
        slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(2), Inches(0.5)).text_frame.text = "Fixed Scale"

    add_dual_slide("Overall Summary", 'overall_auto', 'overall_fixed')
    add_dual_slide(f"Latest Data ({latest_date})", 'latest_auto', 'latest_fixed')
    add_dual_slide("Control Chart (AfterBaking)", 'control_auto', 'control_fixed')

    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)
    return ppt_buf

# --- 主介面邏輯 ---

uploaded_files = st.file_uploader("拖曳 Excel 檔案到這裡 (支援多選)", accept_multiple_files=True, type=['xlsx'])

if uploaded_files:
    if st.button("🚀 開始生成報告"):
        with st.spinner('正在讀取數據並繪製圖表...'):
            df = process_uploaded_files(uploaded_files)
            
            if not df.empty:
                st.success(f"成功讀取 {len(df)} 筆數據！")
                
                # 生成 PPT
                ppt_file = generate_ppt(df)
                
                st.markdown("### ✅ 報告生成完畢！")
                st.download_button(
                    label="📥 點擊下載 PPT 報告",
                    data=ppt_file,
                    file_name="Factory_JMP_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.error("無法解析數據，請檢查 Excel 格式 (是否包含正確的 Sheet 名稱與標題)。")